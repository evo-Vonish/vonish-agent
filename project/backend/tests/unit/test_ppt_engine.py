"""Unit tests for the VonishAgent PPT Artifact Engine.

Exercises the frozen contract in ``ppt_engine`` end to end: the theme/layout
registries, the deterministic SlideIR layout, the python-pptx renderer, the
Pillow preview export, the validator quality gate, the auto-repair loop, the
loose-dict builder, and the ``generate_deck`` orchestrator.

Run from ``F:/Projects/VonishAgent/project/backend``::

    .venv/Scripts/python.exe -m pytest tests/unit/test_ppt_engine.py -q
"""
from __future__ import annotations

import json
import math

import pytest
from pptx import Presentation

from ppt_engine.autorepair import repair_loop
from ppt_engine.builder import build_deck_spec
from ppt_engine.demo_decks import acceptance_deck
from ppt_engine.engine import apply_deck_patch, generate_deck
from ppt_engine.ir import build_deck_ir
from ppt_engine.preview import render_previews
from ppt_engine.registry import get_layout_registry, get_theme_registry
from ppt_engine.renderer import render_pptx
from ppt_engine.schema import (
    BBox,
    Canvas,
    DeckDesignSpec,
    ElementRole,
    ElementType,
    IssueType,
    Severity,
    ShapeStyle,
    ShapeType,
    SlideBackground,
    SlideElement,
    SlideIR,
    TextStyle,
    ValidationReport,
    ValidationSummary,
    ValidatorIssue,
)
from ppt_engine.text_metrics import TEXT_INSET_Y, measure_in_box
from ppt_engine.validator import validate_deck

EXPECTED_THEMES = {
    "tech-dark",
    "academic-white",
    "business-bluegray",
    "vonish-agent",
    "vonish-ocr",
}

EXPECTED_LAYOUTS = {
    "cover-center",
    "toc-simple",
    "chapter-break",
    "three-cards",
    "left-right",
    "timeline",
    "process",
    "architecture",
    "data-chart",
    "quote-center",
    "code-block",
    "summary-bullets",
}

PNG_MAGIC = b"\x89PNG\r\n\x1a\n"


# ---------------------------------------------------------------------------
# Engine-tolerance overflow check
#
# The engine's own definition of "text overflow" lives in the validator
# (PptValidator._check_overflow): it applies a small px/percentage tolerance to
# the raw text-height measurement so a vertically centred single glyph in a
# badge circle is not flagged. We replicate that exact tolerance here so the
# test asserts "no overflow" the way the engine contract defines it, not by a
# stricter raw measurement that would false-positive on legitimate boxes.
# ---------------------------------------------------------------------------
def _engine_overflows(el: SlideElement) -> bool:
    ts = el.text_style
    m = measure_in_box(el.text, ts.fontSize, el.bbox.width, el.bbox.height,
                       line_height=ts.lineHeight, bold=ts.bold)
    over_y = m.text_height > el.bbox.height - 2 * TEXT_INSET_Y + max(6.0, el.bbox.height * 0.04)
    return bool(over_y or m.overflow_x)


def _build_broken_slide(theme) -> SlideIR:
    """A deliberately broken slide that trips overflow, out-of-bounds, off-theme
    colour, tiny font and low-contrast checks — yet whose blocking errors are
    all AUTO-fixable so the repair loop can fully resolve them."""
    bg = SlideBackground(type="solid", color=theme.tokens.background)
    elements = [
        # Overflow: 40pt text far too large for the box (autofit can shrink it).
        SlideElement(
            element_id="overflow", role=ElementRole.BODY, type=ElementType.TEXT,
            bbox=BBox(x=200, y=200, width=420, height=180),
            text="这是一段会溢出的较长中文文字内容示例需要自动缩小字号",
            text_style=TextStyle(fontSize=40, color=theme.tokens.text)),
        # Out-of-bounds (x=-60) AND off-theme magenta fill on one shape.
        SlideElement(
            element_id="oob", role=ElementRole.CARD, type=ElementType.SHAPE,
            bbox=BBox(x=-60, y=100, width=200, height=120),
            shape_type=ShapeType.RECT,
            shape_style=ShapeStyle(fill="#FF00FF")),
        # Font too small (8pt < 12pt minimum).
        SlideElement(
            element_id="tiny", role=ElementRole.BODY, type=ElementType.TEXT,
            bbox=BBox(x=400, y=400, width=300, height=60), text="小字号正文",
            text_style=TextStyle(fontSize=8, color=theme.tokens.text)),
        # Low contrast: near-background dark text on dark background.
        SlideElement(
            element_id="lowc", role=ElementRole.BODY, type=ElementType.TEXT,
            bbox=BBox(x=400, y=560, width=400, height=50),
            text="低对比度文字内容示例展示",
            text_style=TextStyle(fontSize=18, color="#0B0F1B")),
    ]
    return SlideIR(
        deck_id="broken", slide_id="broken-s00", slide_index=0,
        layout_id="summary-bullets", theme_id=theme.theme_id,
        canvas=Canvas(), background=bg, elements=elements)


# ---------------------------------------------------------------------------
# Registries / themes
# ---------------------------------------------------------------------------
def test_registries():
    themes = get_theme_registry()
    layouts = get_layout_registry()

    assert set(themes.list_ids()) == EXPECTED_THEMES
    assert EXPECTED_LAYOUTS.issubset(set(layouts.list_ids()))
    assert len(layouts.list_ids()) == 12

    for theme_id in EXPECTED_THEMES:
        theme = themes.get(theme_id)
        assert theme.theme_id == theme_id
        palette = theme.tokens.palette()
        assert palette, f"{theme_id} palette is empty"
        # palette is lowercased hex
        assert all(c == c.lower() for c in palette)


def test_themes_have_full_tokens():
    required = [
        "background", "surface", "surfaceElevated", "primary", "accent",
        "accentSecondary", "text", "textMuted", "textInverse", "border",
        "borderSubtle", "warning", "success", "error",
    ]
    themes = get_theme_registry()
    for theme_id in EXPECTED_THEMES:
        tokens = themes.get(theme_id).tokens
        for field in required:
            value = getattr(tokens, field)
            assert isinstance(value, str) and value.startswith("#"), (
                f"{theme_id}.{field} is not a hex colour: {value!r}")
            assert len(value.lstrip("#")) in (3, 6)
        assert len(tokens.chart) >= 4, f"{theme_id} chart palette too small"


# ---------------------------------------------------------------------------
# Layout / IR
# ---------------------------------------------------------------------------
def test_build_ir_all_layouts():
    themes = get_theme_registry()
    for theme_id in EXPECTED_THEMES:
        theme = themes.get(theme_id)
        slides = build_deck_ir(acceptance_deck(theme_id), theme)
        assert len(slides) == 12, f"{theme_id} did not build 12 slides"
        for ir in slides:
            elements = ir.all_elements()
            assert len(elements) > 0, (
                f"{theme_id} slide {ir.slide_index} produced no elements")
            for el in elements:
                b = el.bbox
                for coord in (b.x, b.y, b.width, b.height):
                    assert not math.isnan(coord), (
                        f"{theme_id} slide {ir.slide_index} element "
                        f"{el.element_id} has NaN geometry")


def test_no_overflow_on_engine_output():
    """After the engine's autofit/repair, no text element overflows by the
    engine's own overflow definition, and the report has zero TEXT_OVERFLOW."""
    theme = get_theme_registry().get("tech-dark")
    slides = build_deck_ir(acceptance_deck("tech-dark"), theme)
    report, _ = repair_loop(slides, theme, max_rounds=3)

    # The validator (engine's gate) must report no overflow issues.
    overflow_issues = [i for i in report.issues if i.type == IssueType.TEXT_OVERFLOW]
    assert overflow_issues == [], (
        f"engine left {len(overflow_issues)} overflow issue(s)")

    # And measuring every text box with the engine tolerance agrees.
    for ir in slides:
        for el in ir.all_elements():
            if el.type.value == "text" and el.text_style and el.text.strip():
                assert not _engine_overflows(el), (
                    f"slide {ir.slide_index} element {el.element_id} "
                    f"({el.role.value}) overflows at {el.text_style.fontSize}pt: "
                    f"{el.text[:30]!r}")


# ---------------------------------------------------------------------------
# Renderer
# ---------------------------------------------------------------------------
def test_renderer_opens(tmp_path):
    theme = get_theme_registry().get("tech-dark")
    slides = build_deck_ir(acceptance_deck("tech-dark"), theme)
    out_path = tmp_path / "deck.pptx"

    returned = render_pptx(slides, out_path, theme.tokens)
    assert out_path.exists()
    assert returned

    prs = Presentation(str(out_path))
    rendered = list(prs.slides)
    assert len(rendered) == 12
    for slide in rendered:
        assert len(slide.shapes) > 0, "rendered slide has no shapes"


# ---------------------------------------------------------------------------
# Previews
# ---------------------------------------------------------------------------
def test_previews_created(tmp_path):
    theme = get_theme_registry().get("business-bluegray")
    slides = build_deck_ir(acceptance_deck("business-bluegray"), theme)
    out_dir = tmp_path / "previews"

    names = render_previews(slides, theme.tokens, out_dir, scale=1.0)
    assert len(names) == 12

    for name in names:
        png = out_dir / name
        assert png.exists(), f"missing preview {name}"
        data = png.read_bytes()
        assert len(data) > 1024, f"{name} smaller than 1KB ({len(data)} bytes)"
        assert data[:8] == PNG_MAGIC, f"{name} is not a valid PNG"


# ---------------------------------------------------------------------------
# Validator
# ---------------------------------------------------------------------------
def test_validator_finds_issues():
    theme = get_theme_registry().get("tech-dark")
    report = validate_deck([_build_broken_slide(theme)], theme)

    found_types = {issue.type for issue in report.issues}
    required = {
        IssueType.TEXT_OVERFLOW,
        IssueType.OUT_OF_BOUNDS,
        IssueType.COLOR_OUT_OF_THEME,
        IssueType.FONT_TOO_SMALL,
        IssueType.LOW_CONTRAST,
    }
    missing = required - found_types
    assert not missing, f"validator missed issue types: {missing}"
    assert len(found_types) >= 5, (
        f"expected >=5 distinct issue types, got {len(found_types)}: {found_types}")
    # Overflow and out-of-bounds are blocking ERRORs.
    assert report.summary.error_count >= 2


# ---------------------------------------------------------------------------
# Auto-repair
# ---------------------------------------------------------------------------
def test_autorepair_fixes():
    theme = get_theme_registry().get("tech-dark")
    slides = [_build_broken_slide(theme)]

    overflow_before = next(e for e in slides[0].all_elements() if e.element_id == "overflow")
    oob_before = next(e for e in slides[0].all_elements() if e.element_id == "oob")
    font_before = overflow_before.text_style.fontSize
    assert oob_before.bbox.x < 0  # genuinely out of bounds to start

    report, rounds = repair_loop(slides, theme, max_rounds=3)

    # All blocking errors removed (and in this case error_count hits 0).
    blocking = {
        IssueType.TEXT_OVERFLOW, IssueType.ELEMENT_OVERLAP,
        IssueType.OUT_OF_BOUNDS, IssueType.EMPTY_SLIDE,
    }
    blocking_errors = [
        i for i in report.issues
        if i.type in blocking and i.severity == Severity.ERROR
    ]
    assert blocking_errors == [], f"blocking errors survived repair: {blocking_errors}"
    assert report.summary.error_count == 0
    assert report.summary.auto_fixed > 0
    assert report.deliverable is True

    overflow_after = next(e for e in slides[0].all_elements() if e.element_id == "overflow")
    oob_after = next(e for e in slides[0].all_elements() if e.element_id == "oob")

    # Overflow element font shrank.
    assert overflow_after.text_style.fontSize < font_before
    # Out-of-bounds element moved fully inside the canvas.
    assert oob_after.bbox.x >= -0.5
    assert oob_after.bbox.x2 <= slides[0].canvas.width + 0.5
    # Off-theme magenta fill mapped into the theme palette.
    palette = set(theme.tokens.palette())
    assert oob_after.shape_style.fill.lower() in palette
    assert oob_after.shape_style.fill.lower() != "#ff00ff"


# ---------------------------------------------------------------------------
# Builder (loose dicts)
# ---------------------------------------------------------------------------
def test_builder_loose_dicts():
    spec = build_deck_spec(
        title="Loose Deck",
        theme_id="academic-white",
        slides=[
            {"layout": "three-cards", "title": "Cards",
             "cards": ["Alpha", "Beta", "Gamma"]},          # cards as list[str]
            {"layout": "summary-bullets", "title": "Summary",
             "bullets": ["one", "two", "three"]},            # items as list[str]
            {"layout": "left-right", "title": "Compare",
             "left": "left side text", "right": "right side text"},  # left as str
            {"layout": "this-layout-does-not-exist", "title": "Fallback"},
        ],
    )

    assert isinstance(spec, DeckDesignSpec)
    assert spec.theme_id == "academic-white"
    assert spec.title == "Loose Deck"
    assert [s.layout_id for s in spec.slides] == [
        "three-cards", "summary-bullets", "left-right", "summary-bullets",
    ]

    cards = spec.slides[0].content.cards
    assert [c.title for c in cards] == ["Alpha", "Beta", "Gamma"]

    assert spec.slides[1].content.bullets == ["one", "two", "three"]

    left = spec.slides[2].content.left
    assert left is not None and left.body == "left side text"

    # Unknown layout falls back to summary-bullets.
    assert spec.slides[3].layout_id == "summary-bullets"

    # Unknown theme also falls back (to tech-dark).
    fallback = build_deck_spec("t", "no-such-theme",
                               [{"layout": "cover-center", "title": "Hi"}])
    assert fallback.theme_id == "tech-dark"


# ---------------------------------------------------------------------------
# Engine end-to-end
# ---------------------------------------------------------------------------
def test_engine_e2e(tmp_path):
    result = generate_deck(acceptance_deck("business-bluegray"), tmp_path)

    assert result.slide_count == 12
    assert len(result.previews) == 12

    # Every preview file exists on disk (paths are workspace-relative).
    for preview in result.previews:
        assert (tmp_path / preview.path).exists(), f"missing {preview.path}"

    # Sidecar artifacts written.
    assert (tmp_path / result.manifest_path).exists()
    assert (tmp_path / result.deck_spec_path).exists()
    assert (tmp_path / result.slide_ir_path).exists()
    assert (tmp_path / result.pptx_path).exists()

    # SlideIR sidecar is valid JSON with 12 slides.
    ir_data = json.loads((tmp_path / result.slide_ir_path).read_text(encoding="utf-8"))
    assert isinstance(ir_data, list) and len(ir_data) == 12

    assert result.validation.deliverable is True
    assert result.validation.delivery_grade in {"perfect", "good", "acceptable"}


# ---------------------------------------------------------------------------
# Engine must not silently ship an undeliverable deck
# ---------------------------------------------------------------------------
def test_engine_blocks_undeliverable_is_not_silent():
    """A deck with an unresolvable blocking error (an empty slide — AGENT fix
    strategy, never auto-fixed) must surface as not deliverable with a
    populated blocking_issue_types list."""
    theme = get_theme_registry().get("tech-dark")
    bg = SlideBackground(type="solid", color=theme.tokens.background)
    # Only a decorative accent bar — essentially no content.
    empty = SlideIR(
        deck_id="empty", slide_id="empty-s00", slide_index=0,
        layout_id="summary-bullets", theme_id=theme.theme_id,
        canvas=Canvas(), background=bg,
        elements=[SlideElement(
            element_id="bar", role=ElementRole.ACCENTBAR, type=ElementType.SHAPE,
            bbox=BBox(x=100, y=100, width=200, height=6),
            shape_type=ShapeType.RECT,
            shape_style=ShapeStyle(fill=theme.tokens.accent))])

    report, _ = repair_loop([empty], theme, max_rounds=3)

    assert report.deliverable is False
    assert report.delivery_grade == "blocked"
    assert report.blocking_issue_types, "blocking_issue_types must not be empty"
    assert IssueType.EMPTY_SLIDE.value in report.blocking_issue_types

    # The grade logic itself: a hand-built report with a blocking error is
    # blocked + undeliverable; a clean report is deliverable.
    blocking_report = ValidationReport(
        summary=ValidationSummary(total_issues=1, error_count=1),
        issues=[ValidatorIssue(type=IssueType.OUT_OF_BOUNDS, severity=Severity.ERROR)],
        deliverable=False, delivery_grade="blocked",
        blocking_issue_types=[IssueType.OUT_OF_BOUNDS.value])
    assert blocking_report.deliverable is False
    assert blocking_report.blocking_issue_types == ["OUT_OF_BOUNDS"]
    assert len(blocking_report.errors()) == 1

    clean_report = ValidationReport()
    assert clean_report.deliverable is True
    assert clean_report.blocking_issue_types == []


# ---------------------------------------------------------------------------
# Phase-2 element patch: round-trip through the persisted SlideIR
# ---------------------------------------------------------------------------
def test_apply_deck_patch(tmp_path):
    result = generate_deck(acceptance_deck("tech-dark"), tmp_path)
    assert result.validation.deliverable is True

    # locate the cover title element id from the persisted slide meta
    title_box = next(e for e in result.slides_meta[0].elements if e.role == "title")
    assert title_box.text  # has original text

    patched = apply_deck_patch(
        tmp_path, result.pptx_path, slide_index=0,
        operations=[
            {"op": "replace_text", "target": title_box.element_id, "value": "PATCHED 标题"},
            {"op": "update_style", "target": title_box.element_id,
             "changes": {"color": "#22D3EE", "bold": True}},
        ],
        reasoning="unit test patch")

    assert patched.validation.deliverable is True
    assert patched.slide_count == result.slide_count

    # the persisted SlideIR reflects the edit
    ir_data = json.loads((tmp_path / patched.slide_ir_path).read_text(encoding="utf-8"))
    title_el = next(e for e in ir_data[0]["elements"] if e["element_id"] == title_box.element_id)
    assert title_el["text"] == "PATCHED 标题"
    assert title_el["text_style"]["color"] == "#22D3EE"

    # the manifest's slide meta also reflects it, and the pptx/preview exist
    patched_box = next(e for e in patched.slides_meta[0].elements if e.element_id == title_box.element_id)
    assert patched_box.text == "PATCHED 标题"
    assert (tmp_path / patched.pptx_path).exists()
    assert all((tmp_path / p.path).exists() for p in patched.previews)


def test_apply_deck_patch_missing_deck(tmp_path):
    import pytest as _pytest
    with _pytest.raises(FileNotFoundError):
        apply_deck_patch(tmp_path, "outputs/ppt/nope/deck.pptx", 0,
                         [{"op": "replace_text", "target": "x", "value": "y"}])


# ---------------------------------------------------------------------------
# Version history + rollback
# ---------------------------------------------------------------------------
def test_versions_history(tmp_path):
    from ppt_engine.engine import list_deck_versions

    res = generate_deck(acceptance_deck("tech-dark"), tmp_path)
    assert [v.kind for v in res.versions] == ["generate"]
    assert res.versions[0].version_id == "v000"

    tid = next(e for e in res.slides_meta[0].elements if e.role == "title").element_id
    upd = apply_deck_patch(tmp_path, res.pptx_path, 0,
                           [{"op": "replace_text", "target": tid, "value": "v1"}],
                           reasoning="first patch")
    assert [v.kind for v in upd.versions] == ["generate", "patch"]
    # both snapshots exist on disk
    deck_dir = (tmp_path / res.pptx_path).parent
    assert (deck_dir / "versions" / "v000.slideir.json").exists()
    assert (deck_dir / "versions" / "v001.slideir.json").exists()
    assert len(list_deck_versions(tmp_path, res.pptx_path)) == 2


def test_restore_version(tmp_path):
    from ppt_engine.engine import restore_deck_version

    res = generate_deck(acceptance_deck("tech-dark"), tmp_path)
    tid = next(e for e in res.slides_meta[0].elements if e.role == "title")
    original = tid.text

    apply_deck_patch(tmp_path, res.pptx_path, 0,
                     [{"op": "replace_text", "target": tid.element_id, "value": "CHANGED"}],
                     reasoning="change")
    rolled = restore_deck_version(tmp_path, res.pptx_path, "v000")

    # title is back to the original, and a restore version was recorded
    restored_box = next(e for e in rolled.slides_meta[0].elements if e.element_id == tid.element_id)
    assert restored_box.text == original
    assert rolled.versions[-1].kind == "restore"
    assert len(rolled.versions) == 3


# ---------------------------------------------------------------------------
# L2 visual QA
# ---------------------------------------------------------------------------
def test_visual_qa_clean_deck(tmp_path):
    result = generate_deck(acceptance_deck("business-bluegray"), tmp_path, visual_qa=True)
    # findings recorded for every slide & metric (12 slides x 4 metrics)
    assert len(result.visual_findings) == 12 * 4
    # a clean engine deck triggers no image-grounded ERROR issues
    image_errors = [i for i in result.validation.issues
                    if i.type in (IssueType.RENDERED_BLANK, IssueType.RENDER_TEXT_MISSING)]
    assert image_errors == []
    assert result.validation.deliverable is True


def test_visual_qa_catches_invisible_text(tmp_path):
    """Text coloured identically to the slide background renders no ink — the
    L2 layer must catch it even though L1 (geometry) sees a valid element."""
    from ppt_engine.preview import render_previews
    from ppt_engine.visual_qa import run_visual_qa

    theme = get_theme_registry().get("tech-dark")
    bg = theme.tokens.background
    elements = [
        # a visible card so the slide is not globally blank
        SlideElement(element_id="card", role=ElementRole.CARD, type=ElementType.SHAPE,
                     bbox=BBox(x=80, y=80, width=300, height=160),
                     shape_type=ShapeType.RECT,
                     shape_style=ShapeStyle(fill=theme.tokens.surfaceElevated)),
        # invisible text: same colour as the background, on the bare background
        SlideElement(element_id="ghost", role=ElementRole.BODY, type=ElementType.TEXT,
                     bbox=BBox(x=80, y=500, width=700, height=80),
                     text="这段文字与背景同色因此根本看不见",
                     text_style=TextStyle(fontSize=24, color=bg)),
    ]
    ir = SlideIR(deck_id="ghost", slide_id="ghost-s00", slide_index=0,
                 layout_id="custom", theme_id="tech-dark",
                 background=SlideBackground(color=bg), elements=elements)

    names = render_previews([ir], theme.tokens, tmp_path)
    findings, issues = run_visual_qa([ir], [str(tmp_path / names[0])], theme)

    assert any(i.type == IssueType.RENDER_TEXT_MISSING for i in issues)
    # and the per-metric finding marks text_presence as not ok
    tp = next(f for f in findings if f.metric == "text_presence")
    assert tp.ok is False


# ---------------------------------------------------------------------------
# End-to-end loop through the actual agent tool executor
# ---------------------------------------------------------------------------
@pytest.mark.anyio
async def test_tool_loop_generate_patch_revert(tmp_path, monkeypatch):
    """The full Workbench loop, driven through the real ToolExecutor:
    generate_presentation -> patch_presentation -> revert_presentation,
    plus the reference formatter the agent sees."""
    from agent.tool_executor import ToolCallRequest, get_tool_executor
    from agent.tool_registry import register_default_tools
    from api.chat import _format_references
    from core.config import settings

    register_default_tools()
    monkeypatch.setattr(settings, "workspace_root", str(tmp_path))
    ex = get_tool_executor()
    wid = "loopconv"

    gen = await ex.execute(ToolCallRequest(
        tool_name="generate_presentation", conversation_id=wid, workspace_id=wid,
        arguments={"title": "Loop", "theme_id": "tech-dark", "slides": [
            {"layout": "cover-center", "title": "原标题", "subtitle": "s", "meta": "m"},
            {"layout": "summary-bullets", "title": "要点", "bullets": ["a", "b", "c"]},
        ]}))
    assert gen.success and gen.result["open_artifact"]
    art = gen.result["artifact"]
    deck_path = art["path"]
    assert art["kind"] == "presentation"
    assert len(art["versions"]) == 1

    # the agent would read the selected element from the reference block:
    ref_text = _format_references([{
        "sourceType": "slide-element", "title": "cover title",
        "preview": "原标题",
        "location": {"filePath": deck_path, "slideIndex": 0, "elementId": "el-title-2"},
    }])
    assert "elementId=el-title-2" in ref_text
    assert f"filePath={deck_path}" in ref_text

    # find the real title element id (engine ids are stable but assert generally)
    import json as _json
    man = _json.loads((tmp_path / wid / art["manifestPath"]).read_text(encoding="utf-8"))
    tid = next(e["element_id"] for e in man["slides_meta"][0]["elements"] if e["role"] == "title")

    patched = await ex.execute(ToolCallRequest(
        tool_name="patch_presentation", conversation_id=wid, workspace_id=wid,
        arguments={"deck_path": deck_path, "slide_index": 0, "operations": [
            {"op": "replace_text", "target": tid, "value": "新标题"}]}))
    assert patched.success and patched.result["open_artifact"]
    assert len(patched.result["artifact"]["versions"]) == 2

    reverted = await ex.execute(ToolCallRequest(
        tool_name="revert_presentation", conversation_id=wid, workspace_id=wid,
        arguments={"deck_path": deck_path, "version_id": "v000"}))
    assert reverted.success and reverted.result["open_artifact"]
    man2 = _json.loads((tmp_path / wid / reverted.result["artifact"]["manifestPath"]).read_text(encoding="utf-8"))
    title_now = next(e["text"] for e in man2["slides_meta"][0]["elements"] if e["element_id"] == tid)
    assert title_now == "原标题"
