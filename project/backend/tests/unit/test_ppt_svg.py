"""Unit tests for the Phase-3 EXPERIMENTAL SVG render route.

Exercises ``ppt_engine.svg_renderer`` — the SlideIR -> SVG -> native DrawingML
feasibility experiment. The SVG route is NOT part of the engine main chain;
these tests prove the round-trip produces openable, native (editable) pptx.

Run from ``F:/Projects/VonishAgent/project/backend``::

    .venv/Scripts/python.exe -m pytest tests/unit/test_ppt_svg.py -q
"""
from __future__ import annotations

import xml.etree.ElementTree as ET
from io import BytesIO

from pptx import Presentation

from ppt_engine.demo_decks import acceptance_deck
from ppt_engine.ir import build_deck_ir
from ppt_engine.registry import get_theme_registry
from ppt_engine.svg_renderer import (
    SvgSlideRenderer,
    compare_routes,
    render_deck_via_svg,
)


def _tech_dark_slides():
    theme = get_theme_registry().get("tech-dark")
    return theme, build_deck_ir(acceptance_deck("tech-dark"), theme)


# ---------------------------------------------------------------------------
# slide_to_svg
# ---------------------------------------------------------------------------
def test_slide_to_svg_valid():
    _theme, slides = _tech_dark_slides()
    svg = SvgSlideRenderer().slide_to_svg(slides[0])

    assert isinstance(svg, str)
    assert "<svg" in svg and "</svg>" in svg
    assert "<text" in svg, "expected at least one text node"
    assert "<rect" in svg, "expected at least one rect node (background + shapes)"

    # parses as well-formed XML
    root = ET.fromstring(svg)
    assert root.tag.endswith("svg")
    assert root.get("viewBox") == "0 0 1280 720"


# ---------------------------------------------------------------------------
# svg_to_drawingml
# ---------------------------------------------------------------------------
def test_svg_to_drawingml_opens():
    _theme, slides = _tech_dark_slides()
    renderer = SvgSlideRenderer()
    svg = renderer.slide_to_svg(slides[0])

    data = renderer.svg_to_drawingml(svg)
    assert isinstance(data, bytes) and len(data) > 0

    prs = Presentation(BytesIO(data))
    rendered = list(prs.slides)
    assert len(rendered) >= 1, "svg_to_drawingml produced no slide"
    assert len(list(rendered[0].shapes)) >= 1, "native slide has no shapes"


# ---------------------------------------------------------------------------
# render_deck_via_svg
# ---------------------------------------------------------------------------
def test_render_deck_via_svg(tmp_path):
    _theme, slides = _tech_dark_slides()
    out = tmp_path / "via_svg.pptx"

    returned = render_deck_via_svg(slides, str(out))
    assert returned and out.exists()

    prs = Presentation(str(out))
    assert len(list(prs.slides)) == 12


# ---------------------------------------------------------------------------
# compare_routes (A/B benchmark)
# ---------------------------------------------------------------------------
def test_compare_routes(tmp_path):
    theme, slides = _tech_dark_slides()
    result = compare_routes(slides, theme, str(tmp_path / "compare"))

    assert "direct" in result and "svg" in result
    assert result["direct"]["slide_count"] == 12
    assert result["svg"]["slide_count"] == 12

    # both produced openable, non-empty native pptx with shapes
    for key in ("direct", "svg"):
        path = result[key]["pptx"]
        prs = Presentation(path)
        assert len(list(prs.slides)) == 12
        assert result[key]["shape_count"] >= 1
        assert result[key]["bytes"] > 0

    assert result["svg"]["svg_files"] == 12
    assert isinstance(result["notes"], list) and result["notes"]
