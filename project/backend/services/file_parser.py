"""File parser service for the Agent system.

Parses uploaded files and extracts text content for indexing.
Supports: PDF, DOCX, TXT, CSV, MD, HTML, code files.
"""

from __future__ import annotations

import csv
import io
from pydantic import BaseModel, Field
from pathlib import Path
from typing import Any

from core.logging import get_logger

logger = get_logger(__name__)

# ---------------------------------------------------------------------------
# Data Structures
# ---------------------------------------------------------------------------


class ParseResult(BaseModel):
    """Result of parsing a file."""

    success: bool
    text: str = ""
    metadata: dict[str, Any] = Field(default_factory=dict)
    error: str | None = None
    page_count: int = 0
    word_count: int = 0

    def to_dict(self) -> dict[str, Any]:
        return {
            "success": self.success,
            "text_preview": self.text[:500] if self.text else "",
            "metadata": self.metadata,
            "error": self.error,
            "page_count": self.page_count,
            "word_count": self.word_count,
        }


# ---------------------------------------------------------------------------
# File Parser Service
# ---------------------------------------------------------------------------


class FileParser:
    """Parser for extracting text content from various file types.

    Supported formats:
    - PDF (via PyMuPDF)
    - DOCX (via python-docx)
    - TXT, MD, CSV, HTML, code files (plain text)
    - JSON (formatted)
    """

    def __init__(self) -> None:
        self._max_text_length = 100_000  # characters

    async def parse(self, file_path: str, mime_type: str = "") -> ParseResult:
        """Parse a file and extract text content.

        Args:
            file_path: Path to the file.
            mime_type: MIME type of the file.

        Returns:
            ParseResult with extracted text.
        """
        path = Path(file_path)

        if not path.exists():
            return ParseResult(
                success=False, error=f"File not found: {file_path}"
            )

        # Determine parser by extension or mime type
        ext = path.suffix.lower()

        try:
            if ext == ".pdf" or mime_type == "application/pdf":
                return await self._parse_pdf(path)
            elif ext == ".doc" or mime_type == "application/msword":
                return ParseResult(
                    success=False,
                    error="Legacy .doc parsing is not supported yet. The file was saved to workspace.",
                    metadata={"format": "doc"},
                )
            elif ext == ".docx" or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return await self._parse_docx(path)
            elif ext == ".ppt" or mime_type == "application/vnd.ms-powerpoint":
                return ParseResult(
                    success=False,
                    error="Legacy .ppt parsing is not supported yet. The file was saved to workspace.",
                    metadata={"format": "ppt"},
                )
            elif ext == ".pptx" or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return await self._parse_pptx(path)
            elif ext == ".csv" or mime_type == "text/csv":
                return await self._parse_csv(path)
            elif ext in (".html", ".htm") or mime_type == "text/html":
                return await self._parse_html(path)
            elif ext == ".json" or mime_type == "application/json":
                return await self._parse_json(path)
            else:
                # Plain text (txt, md, code files, etc.)
                return await self._parse_text(path)

        except Exception as e:
            logger.error(f"Parse error for {file_path}: {e}")
            return ParseResult(success=False, error=str(e))

    async def _parse_pdf(self, path: Path) -> ParseResult:
        """Parse PDF file using PyMuPDF."""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(path))
            text_parts = []
            total_pages = len(doc)

            for page_num in range(total_pages):
                page = doc[page_num]
                text_parts.append(f"\n--- Page {page_num + 1} ---\n")
                text_parts.append(page.get_text())

            doc.close()

            full_text = "\n".join(text_parts).strip()
            if not full_text:
                return ParseResult(
                    success=False,
                    error="No extractable text layer detected in PDF.",
                    page_count=total_pages,
                    metadata={"format": "pdf", "pages": total_pages},
                )
            word_count = len(full_text.split())

            return ParseResult(
                success=True,
                text=full_text[: self._max_text_length],
                page_count=total_pages,
                word_count=word_count,
                metadata={"format": "pdf", "pages": total_pages},
            )

        except ImportError:
            logger.warning("PyMuPDF not installed, falling back to basic PDF parsing")
            return ParseResult(
                success=False,
                error="PyMuPDF not available for PDF parsing",
            )

    async def _parse_docx(self, path: Path) -> ParseResult:
        """Parse DOCX file using python-docx."""
        try:
            from docx import Document

            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

            full_text = "\n\n".join(paragraphs)
            word_count = len(full_text.split())

            return ParseResult(
                success=True,
                text=full_text[: self._max_text_length],
                word_count=word_count,
                metadata={"format": "docx", "paragraphs": len(paragraphs)},
            )

        except ImportError:
            return ParseResult(
                success=False,
                error="python-docx not available for DOCX parsing",
            )

    async def _parse_pptx(self, path: Path) -> ParseResult:
        """Parse PPTX file using python-pptx."""
        try:
            from pptx import Presentation

            presentation = Presentation(str(path))
            slide_texts: list[str] = []

            for index, slide in enumerate(presentation.slides, start=1):
                parts: list[str] = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        parts.append(text.strip())
                if parts:
                    slide_texts.append(f"[Slide {index}]\n" + "\n".join(parts))

            full_text = "\n\n".join(slide_texts).strip()
            if not full_text:
                return ParseResult(
                    success=False,
                    error="No extractable slide text detected in PPTX.",
                    page_count=len(presentation.slides),
                    metadata={"format": "pptx", "slides": len(presentation.slides)},
                )

            return ParseResult(
                success=True,
                text=full_text[: self._max_text_length],
                word_count=len(full_text.split()),
                page_count=len(presentation.slides),
                metadata={"format": "pptx", "slides": len(presentation.slides)},
            )

        except ImportError:
            return ParseResult(
                success=False,
                error="python-pptx not available for PPTX parsing",
            )

    async def _parse_csv(self, path: Path) -> ParseResult:
        """Parse CSV file."""
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                rows = list(reader)

            lines = []
            for i, row in enumerate(rows[:1000]):  # Limit to 1000 rows
                lines.append(" | ".join(row))

            full_text = "\n".join(lines)
            row_count = len(rows)

            return ParseResult(
                success=True,
                text=full_text[: self._max_text_length],
                word_count=len(full_text.split()),
                metadata={"format": "csv", "rows": row_count, "columns": len(rows[0]) if rows else 0},
            )

        except Exception as e:
            return ParseResult(success=False, error=f"CSV parse error: {e}")

    async def _parse_html(self, path: Path) -> ParseResult:
        """Parse HTML file using BeautifulSoup."""
        try:
            from bs4 import BeautifulSoup

            with open(path, "r", encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f.read(), "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            text = soup.get_text(separator="\n", strip=True)

            return ParseResult(
                success=True,
                text=text[: self._max_text_length],
                word_count=len(text.split()),
                metadata={"format": "html", "title": soup.title.string if soup.title else ""},
            )

        except ImportError:
            # Fallback to plain text
            return await self._parse_text(path)

    async def _parse_json(self, path: Path) -> ParseResult:
        """Parse JSON file."""
        import json

        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)

            formatted = json.dumps(data, indent=2, ensure_ascii=False)

            return ParseResult(
                success=True,
                text=formatted[: self._max_text_length],
                word_count=len(formatted.split()),
                metadata={"format": "json"},
            )

        except json.JSONDecodeError as e:
            return ParseResult(success=False, error=f"Invalid JSON: {e}")

    async def _parse_text(self, path: Path) -> ParseResult:
        """Parse plain text file."""
        last_error: Exception | None = None
        for encoding in ("utf-8", "utf-8-sig", "gbk", "gb18030"):
            try:
                text = path.read_text(encoding=encoding)

                return ParseResult(
                    success=True,
                    text=text[: self._max_text_length],
                    word_count=len(text.split()),
                    metadata={"format": "text", "extension": path.suffix, "encoding": encoding},
                )

            except UnicodeDecodeError as e:
                last_error = e
                continue
            except Exception as e:
                return ParseResult(success=False, error=f"Text parse error: {e}")

        return ParseResult(success=False, error=f"Text decode error: {last_error}")

    async def parse_from_bytes(
        self, file_name: str, file_data: bytes, mime_type: str = ""
    ) -> ParseResult:
        """Parse file from bytes.

        Args:
            file_name: Original file name.
            file_data: File content as bytes.
            mime_type: MIME type.

        Returns:
            ParseResult with extracted text.
        """
        import tempfile

        ext = Path(file_name).suffix
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name

        try:
            result = await self.parse(tmp_path, mime_type)
            return result
        finally:
            Path(tmp_path).unlink(missing_ok=True)


# ---------------------------------------------------------------------------
# Global Instance
# ---------------------------------------------------------------------------

_file_parser: FileParser | None = None


def get_file_parser() -> FileParser:
    """Get the global file parser instance."""
    global _file_parser
    if _file_parser is None:
        _file_parser = FileParser()
    return _file_parser
