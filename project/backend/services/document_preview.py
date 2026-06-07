"""Structured document preview conversion for the File Workbench.

Converts PDF / DOCX / XLSX / PPTX into structured JSON the frontend renders and
makes selectable (for references). Heavy libraries are imported lazily so they
never slow server start-up.

All functions return a dict with ``success: True`` plus structured content, or
``success: False`` with a structured ``error`` object.
"""

from __future__ import annotations

import base64
from pathlib import Path
from typing import Any


def error(code: str, message: str, recoverable: bool = False, suggested_action: str | None = None) -> dict[str, Any]:
    err: dict[str, Any] = {"code": code, "message": message, "recoverable": recoverable}
    if suggested_action:
        err["suggested_action"] = suggested_action
    return {"success": False, "error": err}


# ---------------------------------------------------------------------------
# PDF (PyMuPDF / fitz)
# ---------------------------------------------------------------------------


def preview_pdf(path: Path, max_block_pages: int = 80) -> dict[str, Any]:
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    try:
        pages: list[dict[str, Any]] = []
        truncated = doc.page_count > max_block_pages
        for i in range(min(doc.page_count, max_block_pages)):
            page = doc.load_page(i)
            rect = page.rect
            blocks: list[dict[str, Any]] = []
            for bi, block in enumerate(page.get_text("blocks")):
                text = (block[4] or "").strip()
                if not text:
                    continue
                blocks.append({
                    "id": f"p{i}b{bi}",
                    "text": text[:2000],
                    "bbox": [round(block[0], 1), round(block[1], 1), round(block[2], 1), round(block[3], 1)],
                })
            pages.append({
                "index": i,
                "width": round(rect.width, 1),
                "height": round(rect.height, 1),
                "blocks": blocks,
            })
        return {
            "success": True,
            "kind": "pdf",
            "pageCount": doc.page_count,
            "blockPages": len(pages),
            "truncated": truncated,
            "pages": pages,
        }
    finally:
        doc.close()


def render_pdf_page(path: Path, page_index: int, scale: float = 1.5) -> dict[str, Any]:
    import fitz

    doc = fitz.open(str(path))
    try:
        if page_index < 0 or page_index >= doc.page_count:
            return error("PAGE_OUT_OF_RANGE", f"Page {page_index} out of range (0..{doc.page_count - 1})")
        page = doc.load_page(page_index)
        scale = max(0.5, min(3.0, scale))
        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        image = base64.b64encode(pix.tobytes("png")).decode("ascii")
        return {
            "success": True,
            "index": page_index,
            "width": pix.width,
            "height": pix.height,
            "image": f"data:image/png;base64,{image}",
        }
    finally:
        doc.close()


# ---------------------------------------------------------------------------
# DOCX (python-docx)
# ---------------------------------------------------------------------------


def preview_docx(path: Path) -> dict[str, Any]:
    import docx
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = docx.Document(str(path))
    blocks: list[dict[str, Any]] = []
    index = 0
    table_index = 0

    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if not text:
                index += 1
                continue
            style = (paragraph.style.name if paragraph.style else "") or ""
            style_lower = style.lower()
            if style_lower.startswith("heading") or style_lower.startswith("title"):
                digits = "".join(ch for ch in style if ch.isdigit())
                blocks.append({
                    "id": f"b{index}",
                    "type": "heading",
                    "level": int(digits) if digits else 1,
                    "text": text,
                })
            elif style_lower.startswith("list") or child.find(qn("w:numPr")) is not None:
                blocks.append({"id": f"b{index}", "type": "list_item", "text": text})
            else:
                blocks.append({"id": f"b{index}", "type": "paragraph", "text": text})
            index += 1
        elif child.tag == qn("w:tbl"):
            table = Table(child, document)
            rows = [[(cell.text or "").strip() for cell in row.cells] for row in table.rows]
            blocks.append({"id": f"t{table_index}", "type": "table", "rows": rows})
            table_index += 1
            index += 1

    return {"success": True, "kind": "docx", "blocks": blocks}


# ---------------------------------------------------------------------------
# XLSX (openpyxl)
# ---------------------------------------------------------------------------


def preview_xlsx(path: Path, max_rows: int = 300, max_cols: int = 40) -> dict[str, Any]:
    import openpyxl

    workbook = openpyxl.load_workbook(str(path), data_only=True)
    try:
        sheets: list[dict[str, Any]] = []
        for worksheet in workbook.worksheets:
            rows: list[list[str]] = []
            truncated = False
            for r, row in enumerate(worksheet.iter_rows(values_only=True)):
                if r >= max_rows:
                    truncated = True
                    break
                cells: list[str] = []
                for c, value in enumerate(row):
                    if c >= max_cols:
                        truncated = True
                        break
                    cells.append("" if value is None else str(value))
                rows.append(cells)
            try:
                merges = [str(rng) for rng in worksheet.merged_cells.ranges]
            except Exception:
                merges = []
            sheets.append({
                "name": worksheet.title,
                "rows": rows,
                "merges": merges,
                "truncated": truncated,
                "maxRow": worksheet.max_row,
                "maxCol": worksheet.max_column,
            })
        return {"success": True, "kind": "xlsx", "sheets": sheets}
    finally:
        workbook.close()


# ---------------------------------------------------------------------------
# PPTX (python-pptx)
# ---------------------------------------------------------------------------


def preview_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(path))
    slide_w = float(prs.slide_width or 9144000)
    slide_h = float(prs.slide_height or 6858000)
    scale = 960.0 / slide_w  # normalise to ~960px wide canvas

    slides: list[dict[str, Any]] = []
    for si, slide in enumerate(prs.slides):
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None
        title = ""
        elements: list[dict[str, Any]] = []
        for shi, shape in enumerate(slide.shapes):
            text = ""
            etype = "shape"
            try:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    etype = "text"
            except Exception:
                text = ""
            try:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    etype = "image"
            except Exception:
                pass
            if title_shape is not None and shape == title_shape:
                etype = "title"
                title = text or title
            left = float(shape.left or 0)
            top = float(shape.top or 0)
            width = float(shape.width or 0)
            height = float(shape.height or 0)
            elements.append({
                "id": f"s{si}e{shi}",
                "type": etype,
                "text": text[:600],
                "bbox": [round(left * scale, 1), round(top * scale, 1), round(width * scale, 1), round(height * scale, 1)],
            })
        if not title:
            for element in elements:
                if element["text"]:
                    title = element["text"][:60]
                    break
        slides.append({
            "index": si,
            "title": title or f"Slide {si + 1}",
            "elements": elements,
            "width": round(slide_w * scale, 1),
            "height": round(slide_h * scale, 1),
        })

    return {
        "success": True,
        "kind": "pptx",
        "slideCount": len(slides),
        "slideWidth": round(slide_w * scale, 1),
        "slideHeight": round(slide_h * scale, 1),
        "slides": slides,
    }
