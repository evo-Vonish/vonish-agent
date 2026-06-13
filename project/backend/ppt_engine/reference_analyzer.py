"""Phase 3: rule-based Reference Deck Analyzer (v1).

Given a user-supplied reference ``.pptx``, extract a deterministic *style
profile* (palette, fonts, title placement, element mix, coarse layout hints)
without any LLM/OCR/network. The profile then drives two engine-facing helpers:

* :meth:`RuleReferenceDeckAnalyzer.suggest_theme` / ``suggest_layouts`` — pick
  the nearest built-in theme and a plausible layout set, and
* :meth:`RuleReferenceDeckAnalyzer.profile_to_theme` — clone that built-in
  ``Theme`` and overlay the extracted palette so ``generate_deck`` can produce a
  brand-new deck *in the reference style* (``generate_deck(..., theme=...)``).

Everything is best-effort and fully defensive: a malformed / missing file yields
a minimal :class:`ReferenceDeckProfile` whose ``notes`` explain what failed —
``build_profile`` and ``analyze`` never raise out to the caller.

``python-pptx`` is imported lazily so importing this module never forces the
dependency.
"""
from __future__ import annotations

from collections import Counter
from typing import Any, Optional

from .interfaces import ReferenceDeckAnalyzer
from .registry import get_layout_registry, get_theme_registry
from .schema import EMU_PER_PX, ReferenceDeckProfile, Theme
from .validator import color_distance

# How many palette colours to keep, and how close two hex colours must be to be
# treated as "the same" colour (Euclidean RGB distance) when de-duping.
_MAX_PALETTE = 8
_DEDUPE_DISTANCE = 24.0
_TITLE_BAND = 1.0 / 3.0   # top / middle / bottom thirds


# ---------------------------------------------------------------------------
# colour helpers
# ---------------------------------------------------------------------------
def _hex_from_rgb(rgb: Any) -> str:
    """``python-pptx`` ``RGBColor`` (or int) -> ``#RRGGBB`` lowercase, or ""."""
    try:
        # RGBColor stringifies to the 6-char uppercase hex (no '#').
        s = str(rgb).strip().lstrip("#")
        if len(s) == 6:
            int(s, 16)  # validate
            return "#" + s.lower()
    except Exception:
        pass
    try:
        val = int(rgb)
        if 0 <= val <= 0xFFFFFF:
            return "#%06x" % val
    except (TypeError, ValueError):
        pass
    return ""


def _is_valid_hex(value: str) -> bool:
    if not isinstance(value, str) or len(value) != 7 or value[0] != "#":
        return False
    try:
        int(value[1:], 16)
        return True
    except ValueError:
        return False


def _dedupe_palette(ranked: list[str]) -> list[str]:
    """Keep colours in frequency order, dropping near-identical duplicates."""
    kept: list[str] = []
    for color in ranked:
        if any(color_distance(color, k) <= _DEDUPE_DISTANCE for k in kept):
            continue
        kept.append(color)
        if len(kept) >= _MAX_PALETTE:
            break
    return kept


# ---------------------------------------------------------------------------
# python-pptx shape readers (all defensive)
# ---------------------------------------------------------------------------
def _shape_fill_hex(shape: Any) -> str:
    """Solid fill colour of a shape as ``#RRGGBB``, or "" when not solid."""
    try:
        from pptx.enum.dml import MSO_FILL_TYPE

        fill = shape.fill
        if fill.type != MSO_FILL_TYPE.SOLID:
            return ""
        return _hex_from_rgb(fill.fore_color.rgb)
    except Exception:
        return ""


def _run_color_hex(run: Any) -> str:
    try:
        from pptx.enum.dml import MSO_THEME_COLOR

        color = run.font.color
        if color is None or color.type is None:
            return ""
        # Only explicit RGB colours are trustworthy; theme-colour indices are not
        # resolvable without the theme part, so we skip them.
        rgb = getattr(color, "rgb", None)
        if rgb is None:
            return ""
        return _hex_from_rgb(rgb)
    except Exception:
        return ""


def _run_font_name(run: Any) -> str:
    try:
        name = run.font.name
        if name and isinstance(name, str):
            return name.strip()
    except Exception:
        pass
    return ""


def _classify_type(shape: Any) -> str:
    """Map a python-pptx shape to text/image/table/chart/shape (parse-aligned)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
    except Exception:
        pass
    try:
        if getattr(shape, "has_chart", False):
            return "chart"
    except Exception:
        pass
    try:
        if getattr(shape, "has_table", False):
            return "table"
    except Exception:
        pass
    try:
        if getattr(shape, "has_text_frame", False) and (shape.text_frame.text or "").strip():
            return "text"
    except Exception:
        pass
    return "shape"


def _is_title_shape(shape: Any) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            idx = shape.placeholder_format.type
            return idx in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except Exception:
        pass
    return False


def _is_text(shape: Any) -> bool:
    try:
        return bool(getattr(shape, "has_text_frame", False)
                    and (shape.text_frame.text or "").strip())
    except Exception:
        return False


def _runs(shape: Any) -> list[Any]:
    out: list[Any] = []
    try:
        if not getattr(shape, "has_text_frame", False):
            return out
        for para in shape.text_frame.paragraphs:
            out.extend(para.runs)
    except Exception:
        pass
    return out


def _iter_shapes(shapes: Any) -> list[Any]:
    """Flatten group shapes one level deep (best-effort, recursive)."""
    flat: list[Any] = []
    for shape in shapes:
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                flat.extend(_iter_shapes(shape.shapes))
                continue
        except Exception:
            pass
        flat.append(shape)
    return flat


def _bbox_px(shape: Any) -> Optional[tuple[float, float, float, float]]:
    try:
        return (
            float(shape.left) / EMU_PER_PX,
            float(shape.top) / EMU_PER_PX,
            float(shape.width) / EMU_PER_PX,
            float(shape.height) / EMU_PER_PX,
        )
    except (TypeError, ValueError):
        return None


# ---------------------------------------------------------------------------
# geometry classifiers
# ---------------------------------------------------------------------------
def _classify_position(bbox: tuple[float, float, float, float],
                       sw: float, sh: float) -> str:
    """Vertical (top/middle/bottom) + horizontal (left/center/right) of centre."""
    x, y, w, h = bbox
    cx = (x + w / 2.0) / sw if sw else 0.5
    cy = (y + h / 2.0) / sh if sh else 0.5
    vert = "top" if cy < _TITLE_BAND else ("bottom" if cy > 2 * _TITLE_BAND else "middle")
    horiz = "left" if cx < _TITLE_BAND else ("right" if cx > 2 * _TITLE_BAND else "center")
    return f"{vert}-{horiz}"


def _aspect_ratio(width_emu: Any, height_emu: Any) -> str:
    try:
        w = float(width_emu)
        h = float(height_emu)
    except (TypeError, ValueError):
        return ""
    if w <= 0 or h <= 0:
        return ""
    ratio = w / h
    if abs(ratio - 16.0 / 9.0) < 0.03:
        return "16:9"
    if abs(ratio - 4.0 / 3.0) < 0.03:
        return "4:3"
    if abs(ratio - 16.0 / 10.0) < 0.03:
        return "16:10"
    px_w = round(w / EMU_PER_PX)
    px_h = round(h / EMU_PER_PX)
    return f"{px_w}x{px_h}"


def _layout_hints(rects: list[tuple[float, float, float, float]],
                  type_counts: dict[str, int],
                  title_positions: list[str],
                  sw: float, sh: float) -> list[str]:
    """Coarse, best-effort layout hints from the geometry / element mix."""
    hints: list[str] = []
    if type_counts.get("chart", 0) > 0 or type_counts.get("table", 0) > 0:
        hints.append("data-heavy")

    # title-centric: a big title sitting in the top/centre band.
    if any(pos.endswith("center") for pos in title_positions):
        hints.append("title-centric")

    # card-grid: >=3 similarly-sized rectangles sharing a horizontal band.
    if sh > 0:
        by_band: dict[int, list[tuple[float, float, float, float]]] = {}
        for r in rects:
            x, y, w, h = r
            cy = (y + h / 2.0) / sh
            by_band.setdefault(int(cy * 4), []).append(r)
        for band in by_band.values():
            if len(band) >= 3:
                widths = sorted(b[2] for b in band)
                if widths[0] > 0 and widths[-1] / max(widths[0], 1.0) <= 1.6:
                    hints.append("card-grid")
                    break

    # two-column: exactly two large blocks split left/right.
    if sw > 0:
        big = [r for r in rects if r[2] >= sw * 0.28 and r[3] >= sh * 0.3]
        left = [r for r in big if (r[0] + r[2] / 2.0) / sw < 0.5]
        right = [r for r in big if (r[0] + r[2] / 2.0) / sw >= 0.5]
        if left and right:
            hints.append("two-column")

    # dedupe, preserve order
    seen: set[str] = set()
    return [h for h in hints if not (h in seen or seen.add(h))]


# ---------------------------------------------------------------------------
# analyzer
# ---------------------------------------------------------------------------
class RuleReferenceDeckAnalyzer(ReferenceDeckAnalyzer):
    """Concrete rule-based reference-deck analyzer (no LLM / OCR / network)."""

    # -- public API ---------------------------------------------------------
    def analyze(self, pptx_path: str) -> dict[str, Any]:
        """Return the extracted profile as a plain dict (never raises)."""
        return self.build_profile(pptx_path).model_dump()

    def build_profile(self, pptx_path: str) -> ReferenceDeckProfile:
        """Open ``pptx_path`` and extract a :class:`ReferenceDeckProfile`."""
        path = str(pptx_path or "")
        try:
            from pptx import Presentation
        except Exception as exc:  # pragma: no cover - import guard
            return ReferenceDeckProfile(
                source_path=path,
                notes=f"python-pptx unavailable: {exc}",
            )

        try:
            prs = Presentation(path)
        except Exception as exc:
            return ReferenceDeckProfile(
                source_path=path,
                notes=f"could not open reference pptx: {exc}",
            )

        try:
            return self._extract(prs, path)
        except Exception as exc:  # pragma: no cover - defensive catch-all
            return ReferenceDeckProfile(
                source_path=path,
                notes=f"analysis failed after open: {exc}",
            )

    # -- extraction ---------------------------------------------------------
    def _extract(self, prs: Any, path: str) -> ReferenceDeckProfile:
        try:
            sw_emu = prs.slide_width
            sh_emu = prs.slide_height
        except Exception:
            sw_emu = sh_emu = 0
        aspect = _aspect_ratio(sw_emu, sh_emu)
        sw_px = (float(sw_emu) / EMU_PER_PX) if sw_emu else 1280.0
        sh_px = (float(sh_emu) / EMU_PER_PX) if sh_emu else 720.0

        fill_colors: Counter[str] = Counter()
        text_colors: Counter[str] = Counter()
        fonts: Counter[str] = Counter()
        type_counts: Counter[str] = Counter()
        title_positions: Counter[str] = Counter()
        rects: list[tuple[float, float, float, float]] = []

        slide_count = 0
        notes_parts: list[str] = []
        for slide in prs.slides:
            slide_count += 1
            try:
                shapes = _iter_shapes(list(slide.shapes))
            except Exception:
                continue

            title_done = False
            for shape in shapes:
                try:
                    etype = _classify_type(shape)
                    type_counts[etype] += 1

                    fill = _shape_fill_hex(shape)
                    if fill:
                        fill_colors[fill] += 1

                    bbox = _bbox_px(shape)
                    if bbox is not None and etype == "shape":
                        rects.append(bbox)

                    if etype == "text":
                        for run in _runs(shape):
                            name = _run_font_name(run)
                            if name:
                                fonts[name] += 1
                            rc = _run_color_hex(run)
                            if rc:
                                text_colors[rc] += 1

                    # title position: explicit title placeholder, else first text.
                    if not title_done and (_is_title_shape(shape) or etype == "text"):
                        if bbox is not None:
                            title_positions[_classify_position(bbox, sw_px, sh_px)] += 1
                            title_done = True
                except Exception:
                    continue

        # palette: text colours and fills, ranked by frequency, deduped.
        merged: Counter[str] = Counter()
        merged.update(text_colors)
        merged.update(fill_colors)
        ranked = [c for c, _ in merged.most_common() if _is_valid_hex(c)]
        palette = _dedupe_palette(ranked)

        font_list = [f for f, _ in fonts.most_common()]
        title_pos_list = [p for p, _ in title_positions.most_common()]
        type_count_dict = dict(type_counts)

        hints = _layout_hints(rects, type_count_dict, title_pos_list, sw_px, sh_px)

        profile = ReferenceDeckProfile(
            source_path=path,
            slide_count=slide_count,
            aspect_ratio=aspect,
            palette=palette,
            fonts=font_list,
            title_positions=title_pos_list,
            element_type_counts=type_count_dict,
            layout_hints=hints,
            notes="",
        )
        profile.suggested_theme_id = self.suggest_theme(profile)
        profile.suggested_layouts = self.suggest_layouts(profile)

        if not palette:
            notes_parts.append("no explicit RGB colours found (theme-colour deck?)")
        if not font_list:
            notes_parts.append("no explicit run fonts found")
        profile.notes = "; ".join(notes_parts)
        return profile

    # -- suggestion helpers -------------------------------------------------
    def suggest_theme(self, profile: ReferenceDeckProfile) -> str:
        """Nearest built-in theme id by summed palette colour distance."""
        registry = get_theme_registry()
        ids = registry.list_ids()
        if not profile.palette:
            # No colour evidence: pick a light/dark default sensibly.
            return ids[0] if ids else ""

        best_id = ids[0] if ids else ""
        best_cost = float("inf")
        for tid in ids:
            try:
                theme_palette = registry.get(tid).tokens.palette()
            except Exception:
                continue
            if not theme_palette:
                continue
            # For each extracted colour, distance to its nearest theme colour.
            cost = 0.0
            for c in profile.palette:
                cost += min(color_distance(c, tc) for tc in theme_palette)
            cost /= max(len(profile.palette), 1)
            if cost < best_cost:
                best_cost = cost
                best_id = tid
        return best_id

    def suggest_layouts(self, profile: ReferenceDeckProfile) -> list[str]:
        """Map hints / element mix to valid built-in layout ids."""
        registry = get_layout_registry()
        valid = set(registry.list_ids())
        out: list[str] = []

        def add(layout_id: str) -> None:
            if layout_id in valid and layout_id not in out:
                out.append(layout_id)

        # Cover always first.
        add("cover-center")
        add("toc-simple")

        hints = set(profile.layout_hints)
        if "title-centric" in hints:
            add("chapter-break")
        if "card-grid" in hints:
            add("three-cards")
        if "two-column" in hints:
            add("left-right")
        if "data-heavy" in hints:
            add("data-chart")

        # element-mix driven extras
        counts = profile.element_type_counts or {}
        if counts.get("chart", 0) > 0:
            add("data-chart")
        if counts.get("image", 0) > 0:
            add("three-cards")
        if counts.get("text", 0) > 0:
            add("summary-bullets")

        # always close with a summary layout
        add("summary-bullets")

        # never return empty: fall back to the registry default ordering
        if not out:
            out = [lid for lid in registry.list_ids()][:3]
        return out

    # -- style transfer -----------------------------------------------------
    def profile_to_theme(self, profile: ReferenceDeckProfile,
                         base_theme_id: str = "") -> Theme:
        """Clone the suggested (or given) built-in theme, overlay the palette.

        Text / contrast tokens are kept from the base for safety; only the
        decorative tokens (primary / accent / accentSecondary / background /
        surface) are overridden, and only when the extracted palette gives a
        confident colour. Returns a valid :class:`Theme` with a fresh id so
        ``generate_deck(..., theme=...)`` can render in the reference style.
        """
        registry = get_theme_registry()
        base_id = base_theme_id or profile.suggested_theme_id or self.suggest_theme(profile)
        base = registry.get(base_id)
        clone = base.model_copy(deep=True)
        clone.theme_id = f"ref-{base.theme_id}"
        clone.name = f"Reference style ({base.name})"
        clone.family = "brand"
        clone.description = (
            f"Derived from reference deck '{profile.source_path}' over base "
            f"'{base.theme_id}'."
        )

        palette = [c for c in profile.palette if _is_valid_hex(c)]
        if palette:
            bg = clone.tokens.background
            mode = clone.mode
            # Order extracted colours by distance from the background so the
            # brightest/most-saturated divergent colours become accents.
            ordered = sorted(palette, key=lambda c: -color_distance(c, bg))

            # primary / accent / accentSecondary: the most divergent colours.
            if len(ordered) >= 1:
                clone.tokens.primary = ordered[0]
            if len(ordered) >= 2:
                clone.tokens.accent = ordered[1]
            if len(ordered) >= 3:
                clone.tokens.accentSecondary = ordered[2]

            # background / surface: choose colours matching the theme's mode
            # (darkest for dark themes, lightest for light themes) — only if the
            # reference clearly supplies such a colour.
            def _luma(c: str) -> float:
                h = c.lstrip("#")
                r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                return 0.2126 * r + 0.7152 * g + 0.0722 * b

            by_luma = sorted(palette, key=_luma)
            if mode == "dark":
                cand_bg = by_luma[0]
                if _luma(cand_bg) < 90:
                    clone.tokens.background = cand_bg
            else:
                cand_bg = by_luma[-1]
                if _luma(cand_bg) > 200:
                    clone.tokens.background = cand_bg

            # keep chart palette aligned with the new accents when we have enough
            if len(ordered) >= 3:
                clone.tokens.chart = (ordered[:3] + list(base.tokens.chart))[:6]

        return clone


__all__ = ["RuleReferenceDeckAnalyzer"]
