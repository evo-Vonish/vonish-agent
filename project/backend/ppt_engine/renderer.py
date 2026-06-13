"""PptxRenderer — SlideIR -> native .pptx via python-pptx.

The renderer is *dumb on purpose*: it only translates the already-positioned
SlideIR into PowerPoint objects. No layout, no colour decisions, no font
choices happen here — those are fixed upstream by the theme + layout engine.
Every text box, shape and chart is a real, editable OOXML object (not an
image), which keeps the deliverable editable in PowerPoint / WPS.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from .schema import (
    EMU_PER_PX,
    Align,
    ChartContent,
    SlideElement,
    SlideIR,
    ThemeTokens,
    VAlign,
)
from .text_metrics import TEXT_INSET_X, TEXT_INSET_Y

_ALIGN = {Align.LEFT: PP_ALIGN.LEFT, Align.CENTER: PP_ALIGN.CENTER,
          Align.RIGHT: PP_ALIGN.RIGHT, Align.JUSTIFY: PP_ALIGN.JUSTIFY}
_ANCHOR = {VAlign.TOP: MSO_ANCHOR.TOP, VAlign.MIDDLE: MSO_ANCHOR.MIDDLE,
           VAlign.BOTTOM: MSO_ANCHOR.BOTTOM}
_CHART_TYPE = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
}


def _rgb(hex_color: str | None, fallback: str = "#000000") -> RGBColor:
    h = (hex_color or fallback).lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        h = fallback.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _emu(px: float) -> int:
    return int(round(px * EMU_PER_PX))


class PptxRenderer:
    """Render a list of SlideIR into a single .pptx file."""

    def render(self, slides: list[SlideIR], out_path: str | Path,
               tokens: ThemeTokens | None = None) -> str:
        prs = Presentation()
        if slides:
            cv = slides[0].canvas
            prs.slide_width = Emu(_emu(cv.width))
            prs.slide_height = Emu(_emu(cv.height))
        blank = prs.slide_layouts[6]
        for ir in slides:
            slide = prs.slides.add_slide(blank)
            self._paint_background(slide, ir)
            for el in sorted(ir.elements, key=lambda e: e.z_index):
                self._render_element(slide, el)
        out = Path(out_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        return str(out)

    # -- background --------------------------------------------------------
    def _paint_background(self, slide, ir: SlideIR) -> None:
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(ir.background.color)
        except Exception:
            # Fall back to a full-bleed rectangle if background fill is rejected.
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
                                          Emu(_emu(ir.canvas.width)), Emu(_emu(ir.canvas.height)))
            rect.fill.solid()
            rect.fill.fore_color.rgb = _rgb(ir.background.color)
            rect.line.fill.background()

    # -- elements ----------------------------------------------------------
    def _render_element(self, slide, el: SlideElement) -> None:
        t = el.type.value
        if t == "group":
            for child in sorted(el.children, key=lambda e: e.z_index):
                self._render_element(slide, child)
            return
        if t == "text":
            self._render_text(slide, el)
        elif t == "shape" or t == "line":
            self._render_shape(slide, el)
        elif t == "chart":
            self._render_chart(slide, el)
        # images intentionally unsupported in Phase 1 (no binary assets yet)

    def _render_shape(self, slide, el: SlideElement) -> None:
        b = el.bbox
        st = el.shape_style
        mso = MSO_SHAPE.ROUNDED_RECTANGLE
        if el.shape_type and el.shape_type.value in ("ellipse",):
            mso = MSO_SHAPE.OVAL
        elif el.shape_type and el.shape_type.value in ("rect", "line"):
            mso = MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(mso, Emu(_emu(b.x)), Emu(_emu(b.y)),
                                       Emu(_emu(b.width)), Emu(_emu(b.height)))
        shape.shadow.inherit = False
        if st and st.radius and mso == MSO_SHAPE.ROUNDED_RECTANGLE:
            try:
                frac = max(0.0, min(0.5, st.radius / max(1.0, min(b.width, b.height))))
                shape.adjustments[0] = frac
            except Exception:
                pass
        if st and st.fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(st.fill)
        else:
            shape.fill.background()
        if st and st.stroke and st.strokeWidth > 0:
            shape.line.color.rgb = _rgb(st.stroke)
            shape.line.width = Pt(st.strokeWidth)
        else:
            shape.line.fill.background()
        if el.text:
            self._fill_text_frame(shape.text_frame, el)

    def _render_text(self, slide, el: SlideElement) -> None:
        b = el.bbox
        box = slide.shapes.add_textbox(Emu(_emu(b.x)), Emu(_emu(b.y)),
                                       Emu(_emu(b.width)), Emu(_emu(b.height)))
        self._fill_text_frame(box.text_frame, el)

    def _fill_text_frame(self, tf, el: SlideElement) -> None:
        ts = el.text_style
        tf.word_wrap = True
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        tf.margin_left = Emu(_emu(TEXT_INSET_X))
        tf.margin_right = Emu(_emu(TEXT_INSET_X))
        tf.margin_top = Emu(_emu(TEXT_INSET_Y))
        tf.margin_bottom = Emu(_emu(TEXT_INSET_Y))
        if ts:
            tf.vertical_anchor = _ANCHOR.get(ts.valign, MSO_ANCHOR.TOP)
        paragraphs = (el.text or "").split("\n")
        for i, line in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if ts:
                p.alignment = _ALIGN.get(ts.align, PP_ALIGN.LEFT)
                try:
                    p.line_spacing = ts.lineHeight
                except Exception:
                    pass
            run = p.add_run()
            run.text = line
            if ts:
                run.font.size = Pt(ts.fontSize)
                run.font.name = ts.fontFamily
                run.font.bold = ts.bold
                run.font.italic = ts.italic
                run.font.color.rgb = _rgb(ts.color)

    def _render_chart(self, slide, el: SlideElement) -> None:
        chart: ChartContent = el.chart or ChartContent()
        b = el.bbox
        data = CategoryChartData()
        cats = chart.categories or [f"#{i+1}" for i in range(
            max((len(s.values) for s in chart.series), default=0))]
        data.categories = cats
        if chart.series:
            for s in chart.series:
                vals = list(s.values) + [0] * (len(cats) - len(s.values))
                data.add_series(s.name or "series", vals[:len(cats)])
        else:
            data.add_series("series", [0] * len(cats))
        ctype = _CHART_TYPE.get(chart.type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        try:
            gf = slide.shapes.add_chart(ctype, Emu(_emu(b.x)), Emu(_emu(b.y)),
                                        Emu(_emu(b.width)), Emu(_emu(b.height)), data)
            ch = gf.chart
            ch.has_title = False
            if len(chart.series) > 1:
                ch.has_legend = True
                ch.legend.position = XL_LEGEND_POSITION.BOTTOM
                ch.legend.include_in_layout = False
        except Exception:
            # Charts can fail on exotic inputs; never abort the whole deck.
            box = slide.shapes.add_textbox(Emu(_emu(b.x)), Emu(_emu(b.y)),
                                           Emu(_emu(b.width)), Emu(_emu(b.height)))
            box.text_frame.text = "[chart]"


def render_pptx(slides: list[SlideIR], out_path: str | Path,
                tokens: ThemeTokens | None = None) -> str:
    return PptxRenderer().render(slides, out_path, tokens)
