"""Parse an existing .pptx back into engine types (Phase 2 reverse path).

The renderer turns a ``SlideIR`` into native OOXML. This module does the
inverse just enough for the workbench: open a .pptx with python-pptx, walk its
shapes, and re-express each one as an :class:`ElementBox` / element-tree node so
the workbench can show an element overlay and so a round-trip is possible.

Geometry is converted from EMU back to canvas pixels by dividing by
``EMU_PER_PX`` (the same constant the renderer multiplies by). All work is
best-effort and fully defensive: a malformed shape is skipped, never fatal.

``python-pptx`` is imported lazily inside the functions so importing this module
never forces the dependency.
"""
from __future__ import annotations

from typing import Any, Optional

from .schema import EMU_PER_PX, ElementBox, SlideMeta

_TEXT_LIMIT = 300


def _px(emu: Any) -> float:
    """EMU -> canvas pixels, rounded to 2dp. Tolerates None/garbage."""
    try:
        return round(float(emu) / EMU_PER_PX, 2)
    except (TypeError, ValueError):
        return 0.0


def _shape_text(shape: Any) -> str:
    """Best-effort text extraction, capped at ``_TEXT_LIMIT`` chars."""
    try:
        if getattr(shape, "has_text_frame", False):
            text = (shape.text_frame.text or "").strip()
            return text[:_TEXT_LIMIT]
    except Exception:
        pass
    return ""


def _is_group(shape: Any) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        return shape.shape_type == MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def _classify_type(shape: Any) -> str:
    """Map a python-pptx shape to one of text/shape/image/chart/table."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return "image"
    except Exception:
        st = None
    try:
        if getattr(shape, "has_chart", False):
            return "chart"
    except Exception:
        pass
    try:
        if getattr(shape, "has_table", False):
            return "table"
    except Exception:
        pass
    try:
        if getattr(shape, "has_text_frame", False) and (shape.text_frame.text or "").strip():
            return "text"
    except Exception:
        pass
    return "shape"


def _classify_role(shape: Any, etype: str, is_first_text: bool) -> str:
    """Infer a coarse role; title placeholders win, else first text -> title."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            idx = shape.placeholder_format.type
            if idx in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return "title"
            if idx == PP_PLACEHOLDER.SUBTITLE:
                return "subtitle"
    except Exception:
        pass
    if etype == "image":
        return "image"
    if etype == "chart":
        return "chart"
    if etype == "table":
        return "table"
    if etype == "text":
        return "title" if is_first_text else "body"
    return "decoration"


def _bbox_px(shape: Any) -> list[float]:
    return [_px(shape.left), _px(shape.top), _px(shape.width), _px(shape.height)]


def parse_pptx_to_slide_metas(path: str) -> list[SlideMeta]:
    """Open a .pptx and return one :class:`SlideMeta` per slide.

    Each shape becomes an :class:`ElementBox` with id ``s{si}e{shi}``, a
    best-effort role/type, a px bbox and up to 300 chars of text. Never raises:
    on any failure an empty list is returned.
    """
    try:
        from pptx import Presentation
    except Exception:
        return []

    try:
        prs = Presentation(path)
    except Exception:
        return []

    metas: list[SlideMeta] = []
    for si, slide in enumerate(prs.slides):
        boxes: list[ElementBox] = []
        title = ""
        seen_text = False
        for shi, shape in enumerate(slide.shapes):
            try:
                etype = _classify_type(shape)
                text = _shape_text(shape)
                role = _classify_role(shape, etype, is_first_text=(etype == "text" and not seen_text))
                if etype == "text" and not seen_text:
                    seen_text = True
                    if not title and text:
                        title = text
                if role == "title" and not title and text:
                    title = text
                boxes.append(ElementBox(
                    element_id=f"s{si}e{shi}",
                    role=role,
                    type=etype,
                    bbox=_bbox_px(shape),
                    text=text,
                ))
            except Exception:
                # Per-shape robustness: skip anything we cannot read.
                continue
        metas.append(SlideMeta(
            slide_id=f"parsed-s{si:02d}",
            slide_index=si,
            layout_id="parsed",
            title=title,
            preview="",
            elements=boxes,
        ))
    return metas


def _shape_to_node(shape: Any, eid: str, is_first_text: bool) -> Optional[dict]:
    """One element-tree node; recurse into group shapes as ``children``."""
    try:
        if _is_group(shape):
            children: list[dict] = []
            seen_text = False
            for ci, sub in enumerate(shape.shapes):
                node = _shape_to_node(sub, f"{eid}c{ci}",
                                      is_first_text=False if seen_text else _is_text(sub))
                if node is None:
                    continue
                if node["type"] == "text":
                    seen_text = True
                children.append(node)
            return {
                "element_id": eid,
                "role": "group",
                "type": "group",
                "bbox": _bbox_px(shape),
                "text": "",
                "z_index": 10,
                "children": children,
            }
        etype = _classify_type(shape)
        text = _shape_text(shape)
        role = _classify_role(shape, etype, is_first_text=is_first_text)
        return {
            "element_id": eid,
            "role": role,
            "type": etype,
            "bbox": _bbox_px(shape),
            "text": text,
            "z_index": 10,
        }
    except Exception:
        return None


def _is_text(shape: Any) -> bool:
    try:
        return bool(getattr(shape, "has_text_frame", False) and (shape.text_frame.text or "").strip())
    except Exception:
        return False


def parse_pptx_to_element_tree(path: str) -> dict:
    """Parse a .pptx into a nested element tree dict for the workbench.

    Returns ``{"success": True, "slide_count": n, "slides": [...]}`` where each
    slide holds ``elements`` with recursive ``children`` for group shapes. On
    any failure returns ``{"success": False, "error": str}``.
    """
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - import guard
        return {"success": False, "error": f"python-pptx unavailable: {exc}"}

    try:
        prs = Presentation(path)
    except Exception as exc:
        return {"success": False, "error": f"could not open {path}: {exc}"}

    try:
        slides_out: list[dict] = []
        for si, slide in enumerate(prs.slides):
            elements: list[dict] = []
            title = ""
            seen_text = False
            for shi, shape in enumerate(slide.shapes):
                node = _shape_to_node(shape, f"s{si}e{shi}",
                                      is_first_text=(not seen_text and _is_text(shape)))
                if node is None:
                    continue
                if node["type"] == "text" and not seen_text:
                    seen_text = True
                    if not title and node["text"]:
                        title = node["text"]
                if node["role"] == "title" and not title and node["text"]:
                    title = node["text"]
                elements.append(node)
            slides_out.append({
                "slide_index": si,
                "title": title,
                "elements": elements,
            })
        return {"success": True, "slide_count": len(slides_out), "slides": slides_out}
    except Exception as exc:
        return {"success": False, "error": f"parse failed: {exc}"}
