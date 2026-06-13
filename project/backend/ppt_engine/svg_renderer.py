"""Phase-3 EXPERIMENTAL SVG intermediate route — SVG -> native DrawingML.

This module is a *feasibility experiment* for an alternative render path:

    SlideIR  --slide_to_svg-->  standalone SVG  --svg_to_drawingml-->  .pptx

The hypothesis: a visual middle layer (SVG) is easier to reason about, diff and
even hand-author / LLM-author than DrawingML, while still compiling down to
NATIVE, editable PowerPoint shapes (not a rasterized screenshot). This file
proves the round-trip is feasible with the standard library + python-pptx.

NOT WIRED INTO THE MAIN CHAIN. The production renderer is
``ppt_engine.renderer.render_pptx`` (SlideIR -> DrawingML directly). Nothing in
the engine imports this module; it exists for the research report and the
``compare_routes`` A/B benchmark only.

Known limitations of the SVG route (documented honestly):
  * No gradients — backgrounds and fills collapse to a single solid colour.
  * No shadows / blur / opacity — ``ShapeStyle.shadow`` and ``opacity`` are lost.
  * Text wrapping is APPROXIMATE — a char-count heuristic, not real font metrics
    (the direct route shares the engine's deterministic ``text_metrics`` instead).
  * Charts and tables DEGRADE to placeholder rectangles (no native chart object).
  * Images become a labelled placeholder rectangle (no binary asset embedding).
  * Round-trip SVG parsing only understands the small subset emitted by
    ``slide_to_svg`` (rect / ellipse / line / text / tspan / image-placeholder).
"""
from __future__ import annotations

import io
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from .interfaces import SvgIntermediateRenderer
from .schema import (
    EMU_PER_PX,
    PX_PER_PT,
    Align,
    SlideElement,
    SlideIR,
    ThemeTokens,
)

SVG_NS = "http://www.w3.org/2000/svg"

_ALIGN_ANCHOR = {
    Align.LEFT: "start",
    Align.CENTER: "middle",
    Align.RIGHT: "end",
    Align.JUSTIFY: "start",
}
_ALIGN_PP = {
    "start": PP_ALIGN.LEFT,
    "middle": PP_ALIGN.CENTER,
    "end": PP_ALIGN.RIGHT,
}


# ---------------------------------------------------------------------------
# small helpers
# ---------------------------------------------------------------------------
def _esc(text: str) -> str:
    """Escape a string for inclusion as XML character data / attribute."""
    if not text:
        return ""
    return (text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            .replace('"', "&quot;"))


def _emu(px: float) -> int:
    """Convert SVG user units (px) to EMU."""
    return int(round((px or 0.0) * EMU_PER_PX))


def _norm_hex(value: Optional[str], fallback: str = "#000000") -> str:
    """Normalise a hex colour to ``#rrggbb`` lowercase; fall back if malformed."""
    h = (value or fallback or "#000000").strip()
    if not h.startswith("#"):
        h = "#" + h
    body = h[1:]
    if len(body) == 3:
        body = "".join(c * 2 for c in body)
    if len(body) != 6:
        body = fallback.lstrip("#")
        if len(body) == 3:
            body = "".join(c * 2 for c in body)
    try:
        int(body, 16)
    except ValueError:
        body = "000000"
    return "#" + body.lower()


def _rgb(value: Optional[str], fallback: str = "#000000") -> RGBColor:
    body = _norm_hex(value, fallback).lstrip("#")
    return RGBColor(int(body[0:2], 16), int(body[2:4], 16), int(body[4:6], 16))


def _first_family(stack: Optional[str]) -> str:
    """First concrete family name from a CSS-style font stack."""
    if not stack:
        return "Arial"
    first = stack.split(",")[0].strip().strip('"').strip("'").strip()
    return first or "Arial"


def _wrap_by_chars(text: str, box_width_px: float, font_px: float) -> list[str]:
    """Cheap char-count line wrap (approximate; ignores real glyph metrics).

    Honours explicit newlines. CJK glyphs are ~1 em wide, latin ~0.55 em, so we
    estimate an average advance and split greedily. This is intentionally crude
    — the production route uses ``text_metrics`` instead.
    """
    if not text:
        return []
    avg_advance = max(1.0, font_px * 0.6)
    max_chars = max(1, int(box_width_px / avg_advance))
    out: list[str] = []
    for para in text.split("\n"):
        if para == "":
            out.append("")
            continue
        line = ""
        for ch in para:
            # full-width chars count double toward the budget
            weight = 2 if ord(ch) > 0x2E7F else 1
            projected = sum(2 if ord(c) > 0x2E7F else 1 for c in line) + weight
            if line and projected > max_chars * 1:
                out.append(line)
                line = ch
            else:
                line += ch
        out.append(line)
    return out or [""]


# ---------------------------------------------------------------------------
# Concrete renderer
# ---------------------------------------------------------------------------
class SvgSlideRenderer(SvgIntermediateRenderer):
    """Experimental SVG <-> native-DrawingML round-trip renderer.

    Implements the Phase-3 ``SvgIntermediateRenderer`` ABC. Never raises on a
    bad element — a malformed element is skipped so one bad node cannot abort
    the whole slide (mirrors the defensive style of the direct renderer).
    """

    # -- SlideIR -> SVG ----------------------------------------------------
    def slide_to_svg(self, slide: SlideIR) -> str:
        cv = slide.canvas
        w, h = cv.width, cv.height
        bg = _norm_hex(getattr(slide.background, "color", None) or "#ffffff")
        parts: list[str] = [
            f'<svg xmlns="{SVG_NS}" viewBox="0 0 {w} {h}" '
            f'width="{w}" height="{h}">',
            f'<rect x="0" y="0" width="{w}" height="{h}" fill="{bg}"/>',
        ]
        for el in sorted(slide.elements, key=lambda e: getattr(e, "z_index", 10)):
            try:
                parts.extend(self._element_to_svg(el))
            except Exception:
                # never let a single bad element abort the slide
                continue
        parts.append("</svg>")
        return "\n".join(parts)

    def _element_to_svg(self, el: SlideElement) -> list[str]:
        t = el.type.value
        if t == "group":
            out: list[str] = ["<g>"]
            for child in sorted(el.children, key=lambda e: getattr(e, "z_index", 10)):
                try:
                    out.extend(self._element_to_svg(child))
                except Exception:
                    continue
            out.append("</g>")
            return out
        if t == "text":
            return self._text_to_svg(el)
        if t in ("shape", "line"):
            return self._shape_to_svg(el)
        if t == "image":
            return self._image_to_svg(el)
        if t in ("chart", "table"):
            # degrade to a labelled placeholder (no native chart in SVG route)
            return self._placeholder_to_svg(el, t)
        return []

    def _shape_to_svg(self, el: SlideElement) -> list[str]:
        b = el.bbox
        st = el.shape_style
        fill = _norm_hex(st.fill) if (st and st.fill) else "none"
        stroke = _norm_hex(st.stroke) if (st and st.stroke) else None
        sw = float(st.strokeWidth) if (st and st.strokeWidth) else 0.0
        stroke_attr = (f' stroke="{stroke}" stroke-width="{sw}"'
                       if (stroke and sw > 0) else "")
        stype = el.shape_type.value if el.shape_type else "rect"

        if stype == "ellipse":
            cx, cy = b.x + b.width / 2, b.y + b.height / 2
            node = (f'<ellipse cx="{cx}" cy="{cy}" rx="{b.width / 2}" '
                    f'ry="{b.height / 2}" fill="{fill}"{stroke_attr}/>')
        elif stype == "line":
            # render as a thin rect so width survives the round-trip
            node = (f'<line x1="{b.x}" y1="{b.y + b.height / 2}" '
                    f'x2="{b.x + b.width}" y2="{b.y + b.height / 2}" '
                    f'stroke="{_norm_hex(st.fill if st and st.fill else (st.stroke if st else None))}" '
                    f'stroke-width="{max(1.0, b.height)}"/>')
        else:
            rx = float(st.radius) if (st and st.radius) else 0.0
            rx_attr = f' rx="{rx}"' if rx > 0 else ""
            node = (f'<rect x="{b.x}" y="{b.y}" width="{b.width}" '
                    f'height="{b.height}"{rx_attr} fill="{fill}"{stroke_attr}/>')
        out = [node]
        if el.text:
            out.extend(self._text_to_svg(el))
        return out

    def _text_to_svg(self, el: SlideElement) -> list[str]:
        b = el.bbox
        ts = el.text_style
        if ts is None:
            from .schema import TextStyle
            ts = TextStyle()
        font_px = max(1.0, ts.fontSize * PX_PER_PT)
        family = _first_family(ts.fontFamily)
        color = _norm_hex(ts.color, "#000000")
        weight = "bold" if ts.bold else "normal"
        style = "italic" if ts.italic else "normal"
        anchor = _ALIGN_ANCHOR.get(ts.align, "start")
        inset_x = 10.0
        if anchor == "middle":
            tx = b.x + b.width / 2
        elif anchor == "end":
            tx = b.x + b.width - inset_x
        else:
            tx = b.x + inset_x

        lines = _wrap_by_chars(el.text, max(1.0, b.width - 2 * inset_x), font_px)
        line_h = font_px * (ts.lineHeight or 1.35)
        first_y = b.y + font_px  # baseline of the first line, ~1 em down

        out = [
            f'<text x="{tx}" y="{first_y}" font-size="{font_px}" '
            f'font-family="{_esc(family)}" fill="{color}" '
            f'font-weight="{weight}" font-style="{style}" '
            f'text-anchor="{anchor}">'
        ]
        for i, line in enumerate(lines):
            dy = 0 if i == 0 else line_h
            out.append(f'<tspan x="{tx}" dy="{dy}">{_esc(line)}</tspan>')
        out.append("</text>")
        return out

    def _image_to_svg(self, el: SlideElement) -> list[str]:
        b = el.bbox
        cx, cy = b.x + b.width / 2, b.y + b.height / 2
        return [
            f'<rect x="{b.x}" y="{b.y}" width="{b.width}" height="{b.height}" '
            f'fill="#e5e7eb" stroke="#9ca3af" stroke-width="1"/>',
            f'<text x="{cx}" y="{cy}" font-size="16" font-family="Arial" '
            f'fill="#6b7280" text-anchor="middle">image</text>',
        ]

    def _placeholder_to_svg(self, el: SlideElement, label: str) -> list[str]:
        b = el.bbox
        cx, cy = b.x + b.width / 2, b.y + b.height / 2
        return [
            f'<rect x="{b.x}" y="{b.y}" width="{b.width}" height="{b.height}" '
            f'fill="#eef2f7" stroke="#9ca3af" stroke-width="1"/>',
            f'<text x="{cx}" y="{cy}" font-size="16" font-family="Arial" '
            f'fill="#6b7280" text-anchor="middle">[{_esc(label)}]</text>',
        ]

    # -- SVG -> native DrawingML ------------------------------------------
    def svg_to_drawingml(self, svg: str) -> bytes:
        """Parse one SVG and emit a one-slide .pptx of NATIVE editable shapes.

        Proves the experiment: each SVG primitive becomes a real OOXML object
        (add_shape / add_textbox), not a rasterized image. Returns the saved
        presentation bytes (BytesIO ``.getvalue()``).
        """
        prs = Presentation()
        width, height = self._svg_canvas(svg)
        prs.slide_width = Emu(_emu(width))
        prs.slide_height = Emu(_emu(height))
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        try:
            root = ET.fromstring(svg)
        except ET.ParseError:
            # unparseable SVG -> still return a valid (empty-ish) deck
            return self._save_bytes(prs)
        self._emit_node(slide, root)
        return self._save_bytes(prs)

    def _svg_canvas(self, svg: str) -> tuple[float, float]:
        try:
            root = ET.fromstring(svg)
        except ET.ParseError:
            return (1280.0, 720.0)
        vb = root.get("viewBox")
        if vb:
            try:
                _, _, w, h = (float(v) for v in vb.replace(",", " ").split())
                if w > 0 and h > 0:
                    return (w, h)
            except (ValueError, TypeError):
                pass
        try:
            return (float(root.get("width", 1280)), float(root.get("height", 720)))
        except (ValueError, TypeError):
            return (1280.0, 720.0)

    def _emit_node(self, slide, node) -> None:
        for child in list(node):
            tag = child.tag.split("}")[-1]
            try:
                if tag == "g":
                    self._emit_node(slide, child)
                elif tag == "rect":
                    self._emit_rect(slide, child)
                elif tag == "ellipse":
                    self._emit_ellipse(slide, child)
                elif tag == "line":
                    self._emit_line(slide, child)
                elif tag == "text":
                    self._emit_text(slide, child)
            except Exception:
                # never let one node abort the whole slide
                continue

    def _emit_rect(self, slide, node) -> None:
        x = float(node.get("x", 0)); y = float(node.get("y", 0))
        w = float(node.get("width", 0)); h = float(node.get("height", 0))
        if w <= 0 or h <= 0:
            return
        rx = float(node.get("rx", 0) or 0)
        mso = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(mso, Emu(_emu(x)), Emu(_emu(y)),
                                       Emu(_emu(w)), Emu(_emu(h)))
        shape.shadow.inherit = False
        if mso == MSO_SHAPE.ROUNDED_RECTANGLE:
            try:
                shape.adjustments[0] = max(0.0, min(0.5, rx / max(1.0, min(w, h))))
            except Exception:
                pass
        self._apply_fill_stroke(shape, node)

    def _emit_ellipse(self, slide, node) -> None:
        cx = float(node.get("cx", 0)); cy = float(node.get("cy", 0))
        rx = float(node.get("rx", 0)); ry = float(node.get("ry", 0))
        if rx <= 0 or ry <= 0:
            return
        x, y, w, h = cx - rx, cy - ry, rx * 2, ry * 2
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(_emu(x)), Emu(_emu(y)),
                                       Emu(_emu(w)), Emu(_emu(h)))
        shape.shadow.inherit = False
        self._apply_fill_stroke(shape, node)

    def _emit_line(self, slide, node) -> None:
        x1 = float(node.get("x1", 0)); y1 = float(node.get("y1", 0))
        x2 = float(node.get("x2", 0)); y2 = float(node.get("y2", 0))
        sw = float(node.get("stroke-width", 1) or 1)
        x = min(x1, x2); y = min(y1, y2) - sw / 2
        w = max(1.0, abs(x2 - x1)); h = max(1.0, sw)
        # native thin RECTANGLE stands in for a horizontal rule (editable)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(_emu(x)),
                                       Emu(_emu(y)), Emu(_emu(w)), Emu(_emu(h)))
        shape.shadow.inherit = False
        stroke = node.get("stroke")
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(stroke, "#888888")
        shape.line.fill.background()

    def _apply_fill_stroke(self, shape, node) -> None:
        fill = node.get("fill")
        if fill and fill != "none":
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
        else:
            shape.fill.background()
        stroke = node.get("stroke")
        sw = float(node.get("stroke-width", 0) or 0)
        if stroke and stroke != "none" and sw > 0:
            shape.line.color.rgb = _rgb(stroke)
            shape.line.width = Pt(sw)
        else:
            shape.line.fill.background()

    def _emit_text(self, slide, node) -> None:
        x = float(node.get("x", 0)); y = float(node.get("y", 0))
        font_px = float(node.get("font-size", 18) or 18)
        anchor = node.get("text-anchor", "start")
        color = node.get("fill", "#000000")
        weight = node.get("font-weight", "normal")
        style = node.get("font-style", "normal")
        family = node.get("font-family", "Arial").strip('"').strip("'")

        tspans = [c for c in list(node) if c.tag.split("}")[-1] == "tspan"]
        if tspans:
            lines = ["".join(t.itertext()) for t in tspans]
        else:
            lines = [(node.text or "")]
        lines = [ln for ln in lines if ln is not None]
        if not any(ln.strip() for ln in lines):
            return

        line_h = font_px * 1.35
        n = max(1, len(lines))
        box_h = line_h * n + font_px
        # SVG text y is the baseline of the first line (~1 em below the top)
        top = max(0.0, y - font_px)
        # estimate a generous width so the native box does not re-wrap oddly
        longest = max((len(ln) for ln in lines), default=1)
        est_w = max(120.0, longest * font_px * 0.65)
        if anchor == "middle":
            left = x - est_w / 2
        elif anchor == "end":
            left = x - est_w
        else:
            left = x
        left = max(0.0, left)

        box = slide.shapes.add_textbox(Emu(_emu(left)), Emu(_emu(top)),
                                       Emu(_emu(est_w)), Emu(_emu(box_h)))
        tf = box.text_frame
        tf.word_wrap = True
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        tf.vertical_anchor = MSO_ANCHOR.TOP
        pp_align = _ALIGN_PP.get(anchor, PP_ALIGN.LEFT)
        font_pt = font_px / PX_PER_PT
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = pp_align
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_pt)
            run.font.name = family or "Arial"
            run.font.bold = weight == "bold"
            run.font.italic = style == "italic"
            run.font.color.rgb = _rgb(color, "#000000")

    @staticmethod
    def _save_bytes(prs) -> bytes:
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()


# ---------------------------------------------------------------------------
# Module-level helpers (experimental route)
# ---------------------------------------------------------------------------
def render_deck_to_svgs(slides: list[SlideIR], out_dir: str) -> list[str]:
    """Write one ``slide-NN.svg`` per slide; return the filenames (not paths)."""
    renderer = SvgSlideRenderer()
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    names: list[str] = []
    for i, ir in enumerate(slides):
        try:
            svg = renderer.slide_to_svg(ir)
        except Exception:
            svg = f'<svg xmlns="{SVG_NS}" viewBox="0 0 1280 720"></svg>'
        name = f"slide-{i:02d}.svg"
        (out / name).write_text(svg, encoding="utf-8")
        names.append(name)
    return names


def render_deck_via_svg(slides: list[SlideIR], out_pptx: str,
                        tokens: ThemeTokens | None = None) -> str:
    """Build a multi-slide pptx through the SVG route, then save it.

    For each SlideIR: ``slide_to_svg`` then parse that SVG and add one native
    slide to a shared presentation. Returns the output path.
    """
    renderer = SvgSlideRenderer()
    prs = Presentation()
    if slides:
        cv = slides[0].canvas
        prs.slide_width = Emu(_emu(cv.width))
        prs.slide_height = Emu(_emu(cv.height))
    blank = prs.slide_layouts[6]
    for ir in slides:
        slide = prs.slides.add_slide(blank)
        try:
            svg = renderer.slide_to_svg(ir)
            root = ET.fromstring(svg)
            renderer._emit_node(slide, root)
        except Exception:
            # keep the slide (blank) rather than aborting the deck
            continue
    out = Path(out_pptx)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)


def compare_routes(slides: list[SlideIR], theme, out_dir: str) -> dict:
    """Run BOTH render routes and report a comparison dict.

    direct: ``ppt_engine.renderer.render_pptx`` (SlideIR -> DrawingML).
    svg:    ``render_deck_via_svg`` (SlideIR -> SVG -> DrawingML).
    Each pptx is reopened with python-pptx to count slides and shapes.
    """
    from .renderer import render_pptx  # local import: keep module self-contained

    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    tokens = getattr(theme, "tokens", None) if theme is not None else None

    direct_path = str(out / "direct.pptx")
    svg_path = str(out / "svg_route.pptx")
    svg_dir = str(out / "svg")

    render_pptx(slides, direct_path, tokens)
    svg_files = render_deck_to_svgs(slides, svg_dir)
    render_deck_via_svg(slides, svg_path, tokens)

    def _stats(path: str) -> tuple[int, int, int]:
        prs = Presentation(path)
        deck_slides = list(prs.slides)
        shape_count = sum(len(list(s.shapes)) for s in deck_slides)
        return len(deck_slides), shape_count, len(Path(path).read_bytes())

    d_slides, d_shapes, d_bytes = _stats(direct_path)
    s_slides, s_shapes, s_bytes = _stats(svg_path)

    return {
        "direct": {
            "pptx": direct_path,
            "bytes": d_bytes,
            "slide_count": d_slides,
            "shape_count": d_shapes,
        },
        "svg": {
            "pptx": svg_path,
            "bytes": s_bytes,
            "slide_count": s_slides,
            "shape_count": s_shapes,
            "svg_files": len(svg_files),
        },
        "notes": [
            "EXPERIMENTAL route: not wired into the engine main chain.",
            "SVG route loses gradients, shadows, opacity (solid colours only).",
            "SVG text wrapping is char-count approximate, not font-metric based.",
            "Charts and tables degrade to placeholder rectangles in the SVG route.",
            "Images become labelled placeholder rectangles (no binary embedding).",
            "Both routes emit NATIVE editable OOXML shapes, not rasterized images.",
        ],
    }


__all__ = [
    "SvgSlideRenderer",
    "render_deck_to_svgs",
    "render_deck_via_svg",
    "compare_routes",
]
